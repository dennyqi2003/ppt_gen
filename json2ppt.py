import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_ppt_from_json(json_file, output_file):
    if not os.path.exists(json_file):
        print(f"JSON File {json_file} Not Found")
        return

    try:
        with open(json_file, 'r', encoding='utf-8') as f:
            data_list = json.load(f)
    except json.JSONDecodeError:
        print("Invalid JSON")
        return

    prs = Presentation()
    
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for item in data_list:
        q_id = item.get("id", "")
        description = item.get("description", "")
        full_description = f"{q_id}. {description}"
        
        qa_pairs = item.get("qa_pairs", [])

        if not qa_pairs:
            qa_pairs = [{"question": "（无小题问题）", "answer": "（无答案）"}]

        for qa in qa_pairs:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            left = Inches(0.5)
            top = Inches(0.3)
            width = Inches(12.33)
            height = Inches(3.0)
            
            tb_desc = slide.shapes.add_textbox(left, top, width, height)
            tf_desc = tb_desc.text_frame
            tf_desc.word_wrap = True 
            
            p = tf_desc.paragraphs[0]
            p.text = full_description
            p.font.size = Pt(20) 
            p.font.bold = False

            top_q = Inches(3.5) 
            height_q = Inches(1.5)
            
            tb_q = slide.shapes.add_textbox(left, top_q, width, height_q)
            tf_q = tb_q.text_frame
            tf_q.word_wrap = True
            
            p = tf_q.paragraphs[0]
            question_text = qa.get("question", "")
            p.text = f"【问题】{question_text}"
            p.font.size = Pt(24) 
            p.font.bold = True 
            p.font.color.rgb = RGBColor(0, 51, 153)

            top_a = Inches(5.2)
            height_a = Inches(2.0)
            
            tb_a = slide.shapes.add_textbox(left, top_a, width, height_a)
            tf_a = tb_a.text_frame
            tf_a.word_wrap = True
            
            p = tf_a.paragraphs[0]
            answer_text = qa.get("answer", "")
            p.text = f"【答案】\n{answer_text}"
            p.font.size = Pt(22) 
            p.font.bold = False
            p.font.color.rgb = RGBColor(204, 0, 0)

    # 5. 保存文件
    try:
        prs.save(output_file)
        print(f"Succeeded. PPT saved to {output_file}")
    except PermissionError:
        print(f"Error: Couldn't open {output_file}")

if __name__ == "__main__":
    json_path = "content.json"
    ppt_path = "output.pptx"
    
    create_ppt_from_json(json_path, ppt_path)