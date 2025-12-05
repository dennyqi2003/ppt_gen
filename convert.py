import re
import os
from pptx import Presentation
from pptx.util import Pt, Inches

def read_file_content(filename):
    """
    读取文件内容。
    注意：如果您真的是读取 .docx 格式，建议使用 python-docx 库读取。
    这里保持原逻辑，假设它是文本文件。
    """
    if not os.path.exists(filename):
        print(f"错误: 找不到文件 '{filename}'。请确保文件在当前目录下。")
        return None
        
    try:
        with open(filename, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        print("UTF-8 解码失败，尝试使用 GBK 编码读取...")
        with open(filename, 'r', encoding='gbk') as f:
            return f.read()
    except Exception as e:
        print(f"读取文件时发生未知错误: {e}")
        return None

def parse_exam_text(text):
    """
    解析文本：
    1. 识别题号。
    2. 题号后的所有内容（不管有没有【答案】）都归为一个整体。
    """
    text = text.strip()
    # 正则逻辑保持不变：匹配行首或换行后的数字，后面跟点、顿号、空格或【
    pattern = r"(?:^|\n)\s*(\d+)(?:[\.．\s]|(?=【))"
    parts = re.split(pattern, text)
    
    title = parts[0].strip() if parts[0].strip() else "试卷解析"
    questions = []
    
    # re.split 切割后，parts[0]是标题，后面依次是 [题号, 内容, 题号, 内容...]
    for i in range(1, len(parts), 2):
        q_num = parts[i]
        q_content_raw = parts[i+1]
        
        # --- 修改逻辑开始 ---
        # 不再通过“【答案】”切割，直接获取全部内容
        # 移除首尾空白字符
        full_content = q_content_raw.strip()
        
        questions.append({
            "id": q_num,
            "body": full_content, # 将所有内容都放进 body
            # "answer": ""        # answer 字段不再需要，或者留空
        })
        # --- 修改逻辑结束 ---
        
    return title, questions

def set_font(paragraph, size_pt, is_bold=False):
    """
    辅助函数：设置段落字体
    """
    # 确保段落有内容，否则 runs 可能为空
    if not paragraph.runs:
        paragraph.add_run()

    for run in paragraph.runs:
        run.font.size = Pt(size_pt)
        run.font.bold = is_bold
        run.font.name = 'Microsoft YaHei' # 设置为微软雅黑

def create_ppt(title, questions, output_filename="output_slide.pptx"):
    """
    生成 PPT 文件
    """
    prs = Presentation()
    
    # --- 1. 封面页 ---
    slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    # 如果没有副标题占位符，这里可能会报错，建议加个 try-except 或检查占位符
    try:
        slide.placeholders[1].text = "生成工具"
    except:
        pass

    # --- 设定文本框的位置和大小 ---
    left = top = Inches(0.5)
    width = prs.slide_width - Inches(1)
    height = prs.slide_height - Inches(1)

    # --- 2. 题目页 ---
    for q in questions:
        # 使用 Layout 6 (空白页)
        slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(slide_layout)
        
        # 手动添加一个文本框
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True # 允许自动换行
        
        # --- 添加内容 (题号 + 所有正文) ---
        p_body = tf.add_paragraph()
        p_body.text = f"{q['id']}. {q['body']}"
        set_font(p_body, size_pt=20, is_bold=False) 
        
        # --- 修改逻辑：不再单独处理 answer ---
        # 之前的 answer 处理代码已被移除，因为所有内容都在 body 里了

    # --- 3. 保存 ---
    try:
        prs.save(output_filename)
        print(f"\n成功！已生成文件: {output_filename}")
        print(f"共处理了 {len(questions)} 道题目。")
    except PermissionError:
        print(f"\n错误: 无法写入 '{output_filename}'。请检查文件是否已被打开，关闭后重试。")
    except Exception as e:
        print(f"\n保存时发生错误: {e}")

if __name__ == "__main__":
    # 请确保这里的文件是纯文本内容，或者将 docx 另存为 txt
    # 如果直接读取二进制 docx，read_file_content 可能会读取到乱码
    input_file = "1.docx" 
    output_file = "exam_presentation.pptx"
    
    print(f"正在读取 {input_file} ...")
    content = read_file_content(input_file)
    
    if content:
        # 简单的检查，如果读出来全是乱码(docx头)，可能需要提示用户
        if content.startswith("PK"): 
            print("警告：看起来您正在直接读取 docx 文件，这可能会导致乱码。建议先将 docx 另存为 txt 文件再运行。")

        print("正在解析内容...")
        title_text, parsed_data = parse_exam_text(content)
        
        if not parsed_data:
            print("警告: 未识别到任何题目。请检查文本格式是否包含数字题号（如 1. 或 1、）。")
        else:
            print(f"识别到题目数量: {len(parsed_data)}")
            create_ppt(title_text, parsed_data, output_file)