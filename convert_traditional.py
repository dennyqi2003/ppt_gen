import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt

def generate_basic_ppt(input_file, output_file):

    if not os.path.exists(input_file):
        print(f"Input File {input_file} Not Found")
        return

    with open(input_file, 'r', encoding='utf-8') as f:
        text = f.read()

    pattern = r'(^\s*\d+(?:[.．]|\s))'
    
    parts = re.split(pattern, text, flags=re.MULTILINE)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    count = 0
    for i in range(1, len(parts), 2):
        if i + 1 >= len(parts):
            break
            
        question_no = parts[i].strip()
        content = parts[i+1].strip() 
        
        full_text = f"{question_no} {content}"
        
        slide = prs.slides.add_slide(prs.slide_layouts[6]) 
        
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(12.3)
        height = Inches(6.5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = full_text
        p.font.size = Pt(20)
        
        if len(full_text) > 10:
            count += 1

    try:
        prs.save(output_file)
        print(f"Succeeded. PPT saved to {output_file}")
    except Exception as e:
        print(f"Failed: {e}")

if __name__ == "__main__":
    generate_basic_ppt("input.txt", "output_traditional.pptx")